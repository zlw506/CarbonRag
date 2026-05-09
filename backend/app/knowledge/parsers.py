from __future__ import annotations

import csv
import io
from html.parser import HTMLParser
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
import xlrd

from app.knowledge.schemas import ParsedDocument


class KnowledgeParseError(RuntimeError):
    pass


def parse_document(*, path: Path, mime_type: str) -> ParsedDocument:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return _parsed_document(path=path, mime_type=mime_type, text=_read_text_file(path))
    if suffix in {".html", ".htm"}:
        return _parsed_document(path=path, mime_type=mime_type, text=_read_html_file(path))
    if suffix == ".csv":
        return _parsed_document(path=path, mime_type=mime_type, text=_read_csv_file(path))
    if suffix == ".xlsx":
        return _parsed_document(path=path, mime_type=mime_type, text=_read_xlsx_file(path))
    if suffix == ".xls":
        return _parsed_document(path=path, mime_type=mime_type, text=_read_xls_file(path))
    if suffix == ".docx":
        return _parsed_document(path=path, mime_type=mime_type, text=_read_docx_file(path))
    if suffix == ".pdf":
        return _parsed_document(path=path, mime_type=mime_type, text=_read_pdf_file(path))
    if suffix == ".pptx":
        return _parsed_document(path=path, mime_type=mime_type, text=_read_pptx_file(path))
    if suffix == ".doc":
        raise KnowledgeParseError("Legacy .doc files are not supported. Convert to .docx or extractable PDF first.")
    if suffix == ".ppt":
        raise KnowledgeParseError("Legacy .ppt files are not supported. Convert to .pptx first.")
    raise KnowledgeParseError(f"Unsupported document format: {suffix or 'unknown'}")


def _read_text_file(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gbk"):
        try:
            content = path.read_text(encoding=encoding)
            return _require_text(content)
        except UnicodeDecodeError:
            continue
    raise KnowledgeParseError("Text file encoding is not supported. Use UTF-8 or GBK text.")


def _read_html_file(path: Path) -> str:
    raw_text = _read_text_file(path)
    parser = _ReadableHTMLParser()
    parser.feed(raw_text)
    return _require_text(parser.get_text(), empty_message="HTML file did not contain readable body text.")


def _read_csv_file(path: Path) -> str:
    raw_text = _read_text_file(path)
    reader = csv.reader(io.StringIO(raw_text))
    rows = list(reader)
    if not rows:
        raise KnowledgeParseError("CSV file is empty.")
    headers = rows[0]
    rendered: list[str] = []
    for index, row in enumerate(rows[1:], start=1):
        cells: list[str] = []
        for column, value in zip(headers, row, strict=False):
            if value is None or str(value).strip() == "":
                continue
            cells.append(f"{column}={str(value).strip()}")
        if cells:
            rendered.append(f"[Table csv row {index}]\n" + " | ".join(cells))
    return _require_text("\n\n".join(rendered) or raw_text)


def _read_xlsx_file(path: Path) -> str:
    workbook = load_workbook(filename=path, read_only=True, data_only=True)
    rendered: list[str] = []
    for sheet in workbook.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        headers = [str(cell).strip() if cell is not None else "" for cell in rows[0]]
        for row_index, row in enumerate(rows[1:], start=1):
            pairs: list[str] = []
            for column_name, value in zip(headers, row, strict=False):
                if value is None or str(value).strip() == "":
                    continue
                label = column_name or f"column_{len(pairs) + 1}"
                pairs.append(f"{label}={str(value).strip()}")
            if pairs:
                rendered.append(f"[Sheet {sheet.title} row {row_index}]\n" + " | ".join(pairs))
    return _require_text("\n\n".join(rendered))


def _read_xls_file(path: Path) -> str:
    try:
        workbook = xlrd.open_workbook(path)
    except Exception as exc:  # noqa: BLE001
        raise KnowledgeParseError("Legacy Excel file could not be read.") from exc

    rendered: list[str] = []
    for sheet in workbook.sheets():
        if sheet.nrows == 0:
            continue
        headers = [str(sheet.cell_value(0, col)).strip() for col in range(sheet.ncols)]
        for row_index in range(1, sheet.nrows):
            pairs: list[str] = []
            for col_index in range(sheet.ncols):
                value = sheet.cell_value(row_index, col_index)
                if value is None or str(value).strip() == "":
                    continue
                label = headers[col_index] or f"column_{col_index + 1}"
                pairs.append(f"{label}={str(value).strip()}")
            if pairs:
                rendered.append(f"[Sheet {sheet.name} row {row_index}]\n" + " | ".join(pairs))
    return _require_text("\n\n".join(rendered))


def _read_docx_file(path: Path) -> str:
    document = Document(path)
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in document.tables:
        for row_index, row in enumerate(table.rows, start=1):
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(f"[Table docx row {row_index}]\n" + " | ".join(cells))
    return _require_text("\n\n".join(parts))


def _read_pdf_file(path: Path) -> str:
    try:
        reader = PdfReader(str(path))
    except Exception as exc:  # noqa: BLE001
        raise KnowledgeParseError("PDF file could not be read.") from exc

    pages: list[str] = []
    for page_index, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as exc:  # noqa: BLE001
            raise KnowledgeParseError("PDF text extraction failed.") from exc
        normalized = text.strip()
        if normalized:
            pages.append(f"[Page {page_index}]\n{normalized}")
    return _require_text("\n\n".join(pages), empty_message="PDF did not contain extractable text.")


def _read_pptx_file(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:  # noqa: BLE001
        raise KnowledgeParseError("python-pptx is not installed. Install the optional parser dependency first.") from exc

    try:
        presentation = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        raise KnowledgeParseError("PPTX file could not be read.") from exc

    slides: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                texts.append(text.strip())
        if texts:
            slides.append(f"[Slide {slide_index}]\n" + "\n".join(texts))
    return _require_text("\n\n".join(slides), empty_message="PPTX file did not contain readable text.")


def _parsed_document(*, path: Path, mime_type: str, text: str) -> ParsedDocument:
    return ParsedDocument(text=text, mime_type=mime_type, source_path=str(path))


def _require_text(text: str, *, empty_message: str = "File did not contain readable text.") -> str:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    if not normalized:
        raise KnowledgeParseError(empty_message)
    return normalized


class _ReadableHTMLParser(HTMLParser):
    block_tags = {
        "address",
        "article",
        "aside",
        "blockquote",
        "br",
        "dd",
        "div",
        "dl",
        "dt",
        "figcaption",
        "footer",
        "h1",
        "h2",
        "h3",
        "h4",
        "h5",
        "h6",
        "header",
        "hr",
        "li",
        "main",
        "nav",
        "ol",
        "p",
        "pre",
        "section",
        "table",
        "td",
        "th",
        "tr",
        "ul",
    }
    ignored_tags = {"script", "style", "noscript"}

    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []
        self._ignored_depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:  # noqa: ANN001
        normalized = tag.lower()
        if normalized in self.ignored_tags:
            self._ignored_depth += 1
            return
        if normalized in self.block_tags:
            self._append_break()

    def handle_endtag(self, tag: str) -> None:
        normalized = tag.lower()
        if normalized in self.ignored_tags and self._ignored_depth > 0:
            self._ignored_depth -= 1
            return
        if normalized in self.block_tags:
            self._append_break()

    def handle_data(self, data: str) -> None:
        if self._ignored_depth:
            return
        normalized = " ".join(data.split())
        if normalized:
            self._parts.append(normalized)

    def get_text(self) -> str:
        text = " ".join(self._parts)
        chunks = [chunk.strip() for chunk in text.split("\n") if chunk.strip()]
        if chunks:
            return "\n\n".join(chunks)
        return text.strip()

    def _append_break(self) -> None:
        if self._parts and self._parts[-1] != "\n":
            self._parts.append("\n")
