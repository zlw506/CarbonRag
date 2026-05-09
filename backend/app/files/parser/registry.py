from functools import lru_cache
import re
from pathlib import Path
from typing import Any

from app.files.parser.models import ParsedFileChunkMeta, ParsedFileResult
from app.knowledge.parsers import KnowledgeParseError, parse_document


IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg"}


class FileParserRegistry:
    """Docling-first session file parser with lightweight local fallbacks."""

    def __init__(self, *, docling_provider: Any | None = None) -> None:
        self.docling_provider = docling_provider or _build_docling_provider()

    def parse(self, *, path: Path, mime_type: str) -> ParsedFileResult:
        suffix = path.suffix.lower()
        if suffix in IMAGE_SUFFIXES:
            if self.docling_provider.supports(path, mime_type):
                return self._parse_with_docling(path=path, mime_type=mime_type, ocr_used=True)
            raise KnowledgeParseError("图片或扫描件需要 Docling OCR，可用解析器未启用。")

        if self.docling_provider.supports(path, mime_type):
            try:
                parsed = self._parse_with_docling(path=path, mime_type=mime_type, ocr_used=False)
                if parsed.text.strip():
                    return parsed
            except KnowledgeParseError:
                # Text documents can still use deterministic fallback parsers when Docling is absent or brittle.
                pass

        return self._parse_with_fallback(path=path, mime_type=mime_type)

    def _parse_with_docling(self, *, path: Path, mime_type: str, ocr_used: bool) -> ParsedFileResult:
        document = self.docling_provider.parse(path=path, mime_type=mime_type)
        if not document.text.strip() or document.metadata.get("parse_success") is False:
            raise KnowledgeParseError(str(document.metadata.get("parse_error") or "Docling 未提取到可用文本。"))
        counts = _infer_counts(path)
        return ParsedFileResult(
            text=document.text,
            markdown=document.text,
            parser_name=document.parser_name,
            parser_version=str(document.metadata.get("parser_version") or self.docling_provider.parser_version or ""),
            ocr_used=ocr_used,
            page_count=counts.get("page_count"),
            sheet_count=counts.get("sheet_count"),
            slide_count=counts.get("slide_count"),
            summary=_summarize_text(document.text),
            chunk_metadata=[
                ParsedFileChunkMeta(
                    page_number=block.page,
                    section_title=block.section,
                )
                for block in document.blocks
            ],
            metadata=document.metadata,
        )

    @staticmethod
    def _parse_with_fallback(*, path: Path, mime_type: str) -> ParsedFileResult:
        if path.suffix.lower() == ".html":
            raw_text = path.read_text(encoding="utf-8", errors="ignore")
            text = re.sub(r"<[^>]+>", " ", raw_text)
            text = " ".join(text.split())
            if not text:
                raise KnowledgeParseError("HTML 未提取到可用文本。")
            return ParsedFileResult(
                text=text,
                markdown=text,
                parser_name="carbonrag-html-fallback",
                parser_version="1.0",
                summary=_summarize_text(text),
                metadata={"fallback_parser": "html_tag_strip"},
            )
        if path.suffix.lower() == ".pptx":
            text = _read_pptx_text(path)
            return ParsedFileResult(
                text=text,
                markdown=text,
                parser_name="carbonrag-pptx-fallback",
                parser_version="1.0",
                slide_count=_infer_counts(path).get("slide_count"),
                summary=_summarize_text(text),
                metadata={"fallback_parser": "python-pptx"},
            )
        parsed = parse_document(path=path, mime_type=mime_type)
        counts = _infer_counts(path)
        return ParsedFileResult(
            text=parsed.text,
            markdown=parsed.text,
            parser_name="carbonrag-fallback",
            parser_version="1.0",
            ocr_used=False,
            page_count=counts.get("page_count"),
            sheet_count=counts.get("sheet_count"),
            slide_count=counts.get("slide_count"),
            summary=_summarize_text(parsed.text),
            metadata={"fallback_parser": "app.knowledge.parsers.parse_document"},
        )


def _infer_counts(path: Path) -> dict[str, int | None]:
    suffix = path.suffix.lower()
    counts: dict[str, int | None] = {"page_count": None, "sheet_count": None, "slide_count": None}
    if suffix == ".pdf":
        try:
            from pypdf import PdfReader

            counts["page_count"] = len(PdfReader(str(path)).pages)
        except Exception:
            counts["page_count"] = None
    elif suffix == ".xlsx":
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(path, read_only=True, data_only=True)
            counts["sheet_count"] = len(workbook.sheetnames)
            workbook.close()
        except Exception:
            counts["sheet_count"] = None
    elif suffix == ".pptx":
        try:
            from pptx import Presentation

            counts["slide_count"] = len(Presentation(str(path)).slides)
        except Exception:
            counts["slide_count"] = None
    return counts


def _read_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:  # noqa: BLE001
        raise KnowledgeParseError("PPTX 解析需要 python-pptx，可用解析器未安装。") from exc

    parts: list[str] = []
    presentation = Presentation(str(path))
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(f"幻灯片 {slide_index}：{text.strip()}")
    if not parts:
        raise KnowledgeParseError("PPTX 未提取到可用文本。")
    return "\n\n".join(parts)


def _summarize_text(text: str) -> str:
    normalized = " ".join(text.split())
    return normalized[:240]


def _build_docling_provider() -> Any:
    # Delay this import: app.rag.__init__ pulls retrieval modules that also import app.knowledge.
    from app.rag.parser import DoclingParserProvider

    return DoclingParserProvider()


@lru_cache(maxsize=1)
def get_file_parser_registry() -> FileParserRegistry:
    return FileParserRegistry()
