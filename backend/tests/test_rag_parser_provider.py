import json

import pytest
from docx import Document
from openpyxl import Workbook
from reportlab.pdfgen import canvas

from app.core.config import Settings
from app.rag.parser import (
    DefaultParserProvider,
    DoclingParserProvider,
    MinerUParserProvider,
    ParserRegistry,
    get_default_parser_provider,
)


def test_default_parser_provider_supports_current_file_types() -> None:
    provider = DefaultParserProvider()

    assert provider.supports("sample.md", "text/markdown")
    assert provider.supports("sample.txt", "text/plain")
    assert provider.supports("sample.csv", "text/csv")
    assert provider.supports("sample.xlsx")
    assert provider.supports("sample.xls")
    assert provider.supports("sample.docx")
    assert provider.supports("sample.html", "text/html")
    assert provider.supports("sample.pptx")
    assert provider.supports("sample.pdf", "application/pdf")
    assert provider.supports(content_type="text/plain")
    assert not provider.supports("legacy.doc", "application/msword")


def test_default_parser_provider_parse_returns_parsed_document_with_blocks(tmp_path) -> None:
    source = tmp_path / "sample.md"
    source.write_text("# 企业节能方案\n\n照明优化和空调优化。", encoding="utf-8")
    provider = get_default_parser_provider()

    parsed = provider.parse(source, content_type="text/markdown")

    assert parsed.document_id.startswith("parsed-")
    assert parsed.source_uri == str(source)
    assert parsed.source_type == "local_file"
    assert parsed.title == "sample.md"
    assert parsed.text
    assert parsed.blocks
    assert {block.block_type for block in parsed.blocks} >= {"title", "paragraph"}
    assert parsed.metadata["parser_name"] == "carbonrag-default"
    assert parsed.metadata["parser_version"] == "1.0"
    assert parsed.metadata["source_file"] == str(source)
    assert parsed.metadata["parse_success"] is True
    assert parsed.metadata["parse_error"] is None


def test_default_parser_provider_parses_html_without_script_or_style(tmp_path) -> None:
    source = tmp_path / "policy.html"
    source.write_text(
        """
        <html>
          <head><style>.hidden{}</style><script>ignored()</script></head>
          <body><h1>Policy Title</h1><p>Carbon accounting guidance.</p></body>
        </html>
        """,
        encoding="utf-8",
    )
    provider = DefaultParserProvider()

    parsed = provider.parse(source, content_type="text/html")

    assert parsed.metadata["parse_success"] is True
    assert "Policy Title" in parsed.text
    assert "Carbon accounting guidance" in parsed.text
    assert "ignored" not in parsed.text
    assert parsed.blocks
    assert parsed.metadata["block_count"] == len(parsed.blocks)


def test_default_parser_provider_parses_csv_as_table_blocks(tmp_path) -> None:
    source = tmp_path / "inventory.csv"
    source.write_text("item,value\nElectricity,1200\nGas,300\n", encoding="utf-8")
    provider = DefaultParserProvider()

    parsed = provider.parse(source, content_type="text/csv")

    assert parsed.metadata["parse_success"] is True
    assert parsed.blocks
    assert parsed.blocks[0].block_type == "table"
    assert parsed.blocks[0].metadata["marker_type"] == "table"
    assert parsed.metadata["table_count"] == 2


def test_default_parser_provider_parses_xlsx_with_sheet_metadata(tmp_path) -> None:
    source = tmp_path / "inventory.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Scope1"
    sheet.append(["activity", "amount"])
    sheet.append(["fuel", 42])
    workbook.save(source)
    provider = DefaultParserProvider()

    parsed = provider.parse(source)

    assert parsed.metadata["parse_success"] is True
    assert parsed.blocks
    assert parsed.blocks[0].block_type == "table"
    assert parsed.blocks[0].metadata["marker_type"] == "sheet"
    assert "Scope1" in parsed.blocks[0].section


def test_default_parser_provider_parses_docx_tables(tmp_path) -> None:
    source = tmp_path / "policy.docx"
    document = Document()
    document.add_paragraph("Policy overview")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Clause"
    table.rows[0].cells[1].text = "Requirement"
    document.save(source)
    provider = DefaultParserProvider()

    parsed = provider.parse(source)

    assert parsed.metadata["parse_success"] is True
    assert {block.block_type for block in parsed.blocks} >= {"paragraph", "table"}
    assert any(block.metadata.get("marker_type") == "table" for block in parsed.blocks)


def test_default_parser_provider_parses_pdf_page_blocks(tmp_path) -> None:
    source = tmp_path / "policy.pdf"
    pdf = canvas.Canvas(str(source))
    pdf.drawString(72, 720, "Page one carbon policy")
    pdf.showPage()
    pdf.drawString(72, 720, "Page two accounting guidance")
    pdf.save()
    provider = DefaultParserProvider()

    parsed = provider.parse(source, content_type="application/pdf")

    assert parsed.metadata["parse_success"] is True
    assert parsed.metadata["page_count"] == 2
    assert [block.page for block in parsed.blocks] == [1, 2]
    assert parsed.blocks[0].metadata["marker_type"] == "page"


def test_default_parser_provider_parses_pptx_when_optional_dependency_available(tmp_path) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation

    source = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Carbon roadmap"
    presentation.save(source)
    provider = DefaultParserProvider()

    parsed = provider.parse(source)

    assert parsed.metadata["parse_success"] is True
    assert parsed.blocks
    assert parsed.blocks[0].page == 1
    assert parsed.blocks[0].metadata["marker_type"] == "slide"


def test_default_parser_provider_score_stays_between_zero_and_one(tmp_path) -> None:
    source = tmp_path / "sample.txt"
    source.write_text("企业节能改造\n照明优化\n空调优化", encoding="utf-8")
    provider = DefaultParserProvider()

    parsed = provider.parse(file_path=source, content_type="text/plain")
    score = provider.score(parsed)

    assert 0.0 <= score <= 1.0
    assert parsed.quality_score == score


def test_default_parser_provider_records_parse_error(tmp_path) -> None:
    source = tmp_path / "empty.txt"
    source.write_text("", encoding="utf-8")
    provider = DefaultParserProvider()

    parsed = provider.parse(file_path=source, content_type="text/plain")

    assert parsed.text == ""
    assert parsed.blocks == []
    assert parsed.quality_score == 0.0
    assert parsed.metadata["parse_success"] is False
    assert parsed.metadata["parse_error"]
    assert provider.score(parsed) == 0.0


def test_docling_provider_is_safe_when_docling_is_missing(monkeypatch, tmp_path) -> None:
    source = tmp_path / "sample.pdf"
    source.write_text("not a real pdf", encoding="utf-8")
    monkeypatch.setattr("app.rag.parser.importlib.util.find_spec", lambda name: None)
    provider = DoclingParserProvider()

    assert provider.parser_available is False
    assert provider.supports(source, "application/pdf") is False
    parsed = provider.parse(source, content_type="application/pdf")
    assert parsed.metadata["parser_name"] == "docling"
    assert parsed.metadata["parser_available"] is False
    assert parsed.metadata["parse_success"] is False
    assert parsed.metadata["parse_error"]


def test_parser_registry_uses_default_when_configured(tmp_path) -> None:
    source = tmp_path / "sample.md"
    source.write_text("# 默认解析\n\n当前默认解析行为不变。", encoding="utf-8")
    registry = ParserRegistry(preferred_provider=Settings(rag_parser_provider="default").rag_parser_provider)

    parsed = registry.parse(source, content_type="text/markdown")

    assert parsed.metadata["parser_name"] == "carbonrag-default"
    assert parsed.metadata["parse_success"] is True


def test_parser_registry_falls_back_when_docling_unavailable(monkeypatch, tmp_path) -> None:
    source = tmp_path / "sample.md"
    source.write_text("# Fallback\n\nDocling 不可用时回退默认解析。", encoding="utf-8")
    monkeypatch.setattr("app.rag.parser.importlib.util.find_spec", lambda name: None)
    registry = ParserRegistry(preferred_provider="docling")

    parsed = registry.parse(source, content_type="text/markdown")

    assert parsed.metadata["parser_name"] == "carbonrag-default"
    assert parsed.metadata["parse_success"] is True
    assert parsed.metadata["fallback_from"] == "docling"
    assert parsed.metadata["fallback_reason"] == "docling_unavailable_or_unsupported"
    assert parsed.metadata["parser_chain"] == ["docling:unavailable", "mineru:disabled", "default:success"]


class _FakeDoclingDocument:
    def export_to_markdown(self) -> str:
        return "# Docling 标题\n\nDocling 解析内容。"


class _FakeDoclingResult:
    document = _FakeDoclingDocument()


class _FakeDoclingConverter:
    def convert(self, file_path: str) -> _FakeDoclingResult:
        assert file_path
        return _FakeDoclingResult()


def test_docling_provider_parse_metadata_with_mock_converter(tmp_path) -> None:
    source = tmp_path / "sample.pdf"
    source.write_text("mock pdf", encoding="utf-8")
    provider = DoclingParserProvider(converter=_FakeDoclingConverter())

    parsed = provider.parse(source, content_type="application/pdf")

    assert parsed.metadata["parser_name"] == "docling"
    assert parsed.metadata["parser_available"] is True
    assert "parser_version" in parsed.metadata
    assert parsed.metadata["parse_success"] is True
    assert parsed.metadata["parse_error"] is None
    assert parsed.blocks
    assert parsed.quality_score > 0


def test_mineru_provider_is_safe_when_mineru_is_missing(monkeypatch, tmp_path) -> None:
    source = tmp_path / "sample.pdf"
    source.write_text("not a real pdf", encoding="utf-8")
    monkeypatch.setattr("app.rag.parser.importlib.util.find_spec", lambda name: None)
    provider = MinerUParserProvider(enabled=True)

    assert provider.parser_available is False
    assert provider.supports(source, "application/pdf") is False
    parsed = provider.parse(source, content_type="application/pdf")
    assert parsed.metadata["parser_name"] == "mineru"
    assert parsed.metadata["parser_available"] is False
    assert parsed.metadata["parse_success"] is False
    assert parsed.metadata["fallback_reason"] == "mineru_unavailable_or_unsupported"
    assert parsed.metadata["output_format"] == "markdown"


class _ExplodingMinerUConverter:
    def __init__(self) -> None:
        self.calls = 0

    def convert(self, file_path: str) -> object:
        self.calls += 1
        raise AssertionError(f"MinerU should not be called for {file_path}")


def test_parser_registry_does_not_call_mineru_when_disabled(monkeypatch, tmp_path) -> None:
    source = tmp_path / "sample.md"
    source.write_text("# Fallback\n\nMinerU disabled should not be called.", encoding="utf-8")
    monkeypatch.setattr("app.rag.parser.importlib.util.find_spec", lambda name: None)
    converter = _ExplodingMinerUConverter()
    registry = ParserRegistry(
        preferred_provider="docling",
        enable_mineru=False,
        mineru_provider=MinerUParserProvider(enabled=False, converter=converter),
    )

    parsed = registry.parse(source, content_type="text/markdown")

    assert converter.calls == 0
    assert parsed.metadata["parser_name"] == "carbonrag-default"
    assert parsed.metadata["parse_success"] is True
    assert parsed.metadata["parser_chain"] == ["docling:unavailable", "mineru:disabled", "default:success"]


def test_parser_registry_records_mineru_unavailable_fallback_chain(monkeypatch, tmp_path) -> None:
    source = tmp_path / "sample.md"
    source.write_text("# Chain\n\nDocling and MinerU unavailable should fall back.", encoding="utf-8")
    monkeypatch.setattr("app.rag.parser.importlib.util.find_spec", lambda name: None)
    registry = ParserRegistry(preferred_provider="docling", enable_mineru=True)

    parsed = registry.parse(source, content_type="text/markdown")

    assert parsed.metadata["parser_name"] == "carbonrag-default"
    assert parsed.metadata["parse_success"] is True
    assert parsed.metadata["fallback_from"] == "docling"
    assert parsed.metadata["fallback_reason"] == "docling_unavailable_or_unsupported"
    assert parsed.metadata["parser_chain"] == ["docling:unavailable", "mineru:unavailable", "default:success"]


class _FakeMinerUDocument:
    def export_to_markdown(self) -> str:
        return "# MinerU title\n\nParsed PDF content."


class _FakeMinerUResult:
    document = _FakeMinerUDocument()


class _FakeMinerUConverter:
    def convert(self, file_path: str) -> _FakeMinerUResult:
        assert file_path
        return _FakeMinerUResult()


def test_mineru_provider_parse_metadata_with_mock_converter(tmp_path) -> None:
    source = tmp_path / "complex.pdf"
    source.write_text("mock pdf", encoding="utf-8")
    provider = MinerUParserProvider(enabled=True, converter=_FakeMinerUConverter())

    parsed = provider.parse(source, content_type="application/pdf")

    assert parsed.metadata["parser_name"] == "mineru"
    assert parsed.metadata["parser_enabled"] is True
    assert parsed.metadata["parser_available"] is True
    assert parsed.metadata["parse_success"] is True
    assert parsed.metadata["parse_error"] is None
    assert parsed.metadata["fallback_reason"] is None
    assert parsed.metadata["output_format"] == "markdown"
    assert parsed.blocks
    assert parsed.quality_score > 0


def test_parser_metadata_is_task_record_serializable(monkeypatch, tmp_path) -> None:
    source = tmp_path / "sample.md"
    source.write_text("# Metadata\n\nParser metadata can be stored on task records.", encoding="utf-8")
    monkeypatch.setattr("app.rag.parser.importlib.util.find_spec", lambda name: None)
    registry = ParserRegistry(preferred_provider="docling", enable_mineru=True)

    parsed = registry.parse(source, content_type="text/markdown")
    encoded = json.dumps(parsed.metadata, ensure_ascii=False)

    assert "parser_chain" in encoded
    assert "mineru:unavailable" in encoded
